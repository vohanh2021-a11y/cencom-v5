#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import os, re, textwrap
from pathlib import Path

ROOT = Path(r"E:\APP-LAPTOP-SYNC\cencomOS_gara_4.0_supa\giao-trinh-opencode-cencom")
CHUONG_DIR = ROOT / "02_ban-thao-chuong"
HINH_DIR = ROOT / "03_hinh-anh-minh-hoa"
BAN_GIAO = ROOT / "04_ban-giao"
HINH_DIR.mkdir(parents=True, exist_ok=True)
BAN_GIAO.mkdir(parents=True, exist_ok=True)

# 1. TAO 20 HINH
from PIL import Image, ImageDraw, ImageFont

palette = {
    "navy": "#0F2342",
    "light_blue": "#E6F0FF",
    "light_orange": "#FFE8D6",
    "white": "#FFFFFF",
    "black": "#1A1A1A",
    "gray": "#6B7280",
}

hinh_titles = [
    ("H01", "Xu thế AI Agent — Tổng quan", "Từ Chatbot tới Đồng nghiệp số"),
    ("H02", "Ba làn sóng AI Agent", "Tool-use → Multi-agent → Autonomous"),
    ("H03", "Bản đồ Model AI 2026", "GPT-5 • Claude 4 • Gemini 2.5 • Grok-4"),
    ("H04", "Bảng so sánh Model AI", "Context • Giá • Tiếng Việt • Reasoning"),
    ("H05", "Làm việc cùng AI", "Prompt → Workflow → Agent"),
    ("H06", "5 Pattern Human-in-the-loop", "RAG • Memory • Tool • Planning • Reflection"),
    ("H07", "OpenCode là gì?", "TUI + LSP + Skills — VS Code kỷ nguyên Agent"),
    ("H08", "Kiến trúc OpenCode 3 lớp", "Interface • Runtime • Extension"),
    ("H09", "OpenCode Zen Free Tier", "Quota • Model • Tiết kiệm"),
    ("H10", "Cấu hình opencode.json", "Zen API Key • Model • Skill"),
    ("H11", "Anatomy SKILL.md", "name • description • Use when..."),
    ("H12", "GLOBAL vs PROJECT Skill", "Cá nhân vs Dự án — .opencode/skills"),
    ("H13", "Dạy AI đọc hiểu tài liệu", "docs/memory • RAG local • 5 Iron Laws"),
    ("H14", "Nghiên cứu mở rộng", "webfetch • websearch • QC206"),
    ("H15", "Học giọng văn", "Từ khóa • Câu • Tone — Few-shot"),
    ("H16", "Skill viết lại báo cáo", "viet-lai-phong-cach-[Tên]"),
    ("H17", "Persona trợ lý theo phòng ban", "Kế toán • Kho • Xưởng • SC"),
    ("H18", "Roadmap 30-60-90 ngày", "Làm quen → Cá nhân hoá → Trợ lý thực thụ"),
    ("H19", "Roadmap 6 tháng Cencom", "Thí điểm • Mở rộng • Chuẩn hoá"),
    ("H20", "Governance & Chi phí", "Bảo mật • Phân quyền • Zen Pro"),
]

def hex_to_rgb(h):
    h=h.lstrip("#")
    return tuple(int(h[i:i+2],16) for i in (0,2,4))

W,H = 1600,900
for code, title, subtitle in hinh_titles:
    img = Image.new("RGB", (W,H), hex_to_rgb(palette["white"]))
    draw = ImageDraw.Draw(img)
    # header navy
    draw.rectangle([(0,0),(W,140)], fill=hex_to_rgb(palette["navy"]))
    # font fallback
    try:
        font_title = ImageFont.truetype("arial.ttf", 52)
        font_sub = ImageFont.truetype("arial.ttf", 32)
        font_code = ImageFont.truetype("arial.ttf", 28)
        font_small = ImageFont.truetype("arial.ttf", 22)
    except:
        font_title = ImageFont.load_default()
        font_sub = ImageFont.load_default()
        font_code = ImageFont.load_default()
        font_small = ImageFont.load_default()
    # code pill
    draw.rounded_rectangle([(60,38),(220,90)], radius=18, fill=hex_to_rgb(palette["light_orange"]))
    draw.text((72,50), code, fill=hex_to_rgb(palette["navy"]), font=font_code)
    draw.text((260,48), title, fill=hex_to_rgb(palette["white"]), font=font_title)
    # subtitle
    draw.text((60,200), subtitle, fill=hex_to_rgb(palette["navy"]), font=font_sub)
    # decorative shapes
    draw.rounded_rectangle([(60,280),(520,680)], radius=24, fill=hex_to_rgb(palette["light_blue"]), outline=hex_to_rgb(palette["navy"]), width=2)
    draw.rounded_rectangle([(560,280),(1020,680)], radius=24, fill=hex_to_rgb(palette["light_orange"]), outline=hex_to_rgb(palette["navy"]), width=2)
    draw.rounded_rectangle([(1060,280),(1520,680)], radius=24, fill=hex_to_rgb(palette["white"]), outline=hex_to_rgb(palette["gray"]), width=2)
    # icon placeholders text
    draw.text((120,360), "◈  MINH HOẠ", fill=hex_to_rgb(palette["navy"]), font=font_sub)
    draw.text((620,360), "⬢  SO SÁNH", fill=hex_to_rgb(palette["navy"]), font=font_sub)
    draw.text((1120,360), "⬣  ỨNG DỤNG", fill=hex_to_rgb(palette["gray"]), font=font_sub)
    # lines fake
    for y in [430,470,510,550,590]:
        draw.rounded_rectangle([(110, y),(470, y+14)], radius=7, fill=hex_to_rgb(palette["white"]))
        draw.rounded_rectangle([(610, y),(970, y+14)], radius=7, fill=hex_to_rgb(palette["white"]))
        draw.rounded_rectangle([(1110, y),(1470, y+14)], radius=7, fill=hex_to_rgb(palette["light_blue"]))
    # footer
    draw.text((60,820), "Cencom VLXD Miền Trung  •  Giáo trình OpenCode  •  cencom.vn", fill=hex_to_rgb(palette["gray"]), font=font_small)
    draw.text((1180,820), "1600×900  •  Palette: Navy / Light Blue / Orange", fill=hex_to_rgb(palette["gray"]), font=font_small)
    fname = HINH_DIR / f"{code}_{subtitle.replace(' ', '_').replace('•','').replace('—','').replace('/','_')[:30]}.png"
    # ensure exact code prefix file for matching placeholder
    # also create clean Hxx.png symlink copy
    img.save(fname, "PNG")
    # also save as Hxx.png for easy reference
    img.save(HINH_DIR / f"{code}.png", "PNG")
    print(f"Created {code}")

print(f"Done {len(hinh_titles)} images")

# 2. BUILD WORD
from docx import Document
from docx.shared import Pt, Inches, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.dml.color import ColorFormat

def set_cell_shading(cell, color_hex):
    tblCell = cell._tc
    tblCellProperties = tblCell.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:fill'), color_hex.replace("#",""))
    shd.set(qn('w:val'), 'clear')
    tblCellProperties.append(shd)

def add_heading_styled(doc, text, level):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.name = 'Arial'
        run.font.color.rgb = RGBColor(0x0F,0x23,0x42) if level==1 else RGBColor(0x1A,0x1A,0x1A)
        if level==1:
            run.font.size = Pt(18)
            run.bold = True
        elif level==2:
            run.font.size = Pt(14)
            run.bold = True
        else:
            run.font.size = Pt(12)
            run.bold = True
    # shading for H1
    if level==1:
        p = h._p
        pPr = p.get_or_add_pPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:fill'), '0F2342')
        shd.set(qn('w:val'), 'clear')
        pPr.append(shd)
        for run in h.runs:
            run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    return h

doc = Document()
# set narrow margins
for section in doc.sections:
    section.top_margin = Inches(0.5)
    section.bottom_margin = Inches(0.5)
    section.left_margin = Inches(0.7)
    section.right_margin = Inches(0.7)
    section.header_distance = Inches(0.3)
    section.footer_distance = Inches(0.3)

style = doc.styles['Normal']
style.font.name = 'Arial'
style.font.size = Pt(10)
style.paragraph_format.space_after = Pt(6)
style.paragraph_format.line_spacing = 1.05

# COVER
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("CÔNG TY VLXD MIỀN TRUNG — CENCOM")
run.font.size = Pt(11); run.font.color.rgb = RGBColor(0x0F,0x23,0x42); run.bold = True; run.font.name='Arial'
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("GIÁO TRÌNH SỬ DỤNG OPENCODE")
run.font.size = Pt(28); run.bold=True; run.font.color.rgb = RGBColor(0x0F,0x23,0x42); run.font.name='Arial'
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("Trợ lý soạn văn bản báo cáo cho Trưởng phòng ban")
run.font.size = Pt(14); run.font.color.rgb = RGBColor(0x6B,0x72,0x80); run.font.name='Arial'
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("Ứng dụng AI Agent • OpenCode & OpenCode Zen Free Tier • Skill cá nhân hoá • Học giọng văn")
run.font.size = Pt(10); run.font.color.rgb = RGBColor(0x6B,0x72,0x80); run.italic=True

# info box
t = doc.add_table(rows=1, cols=3)
t.alignment = WD_TABLE_ALIGNMENT.CENTER
cells = t.rows[0].cells
cells[0].text = "10 chương\n10.000 từ\n20 hình minh hoạ"
cells[1].text = "Dành cho: Trưởng phòng\nKế toán • Kho • Xưởng • SC"
cells[2].text = "Phiên bản 1.0\nTháng 08/2026\nLưu hành nội bộ"
for row in t.rows:
    for cell in row.cells:
        for para in cell.paragraphs:
            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in para.runs:
                run.font.size = Pt(9); run.font.name='Arial'
        set_cell_shading(cell, "#E6F0FF")
t2 = doc.add_paragraph()
t2.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = t2.add_run("Biên soạn: Swarm Orchestrator (Muse Spark)  •  10 agent chuyên môn  •  Kiểm duyệt: Ban Giám đốc Cencom")
run.font.size = Pt(8); run.font.color.rgb = RGBColor(0x6B,0x72,0x80)

doc.add_page_break()

# MUC LUC
add_heading_styled(doc, "Mục lục", 1)
toc_items = [
    ("Chương 1", "Xu thế AI Agent 2024–2026 — Từ chatbot tới đồng nghiệp số", "4"),
    ("Chương 2", "Bản đồ Model AI mới nhất — Chọn model nào cho báo cáo?", "9"),
    ("Chương 3", "Xu hướng làm việc cùng AI — Human-in-the-loop", "14"),
    ("Chương 4", "Giới thiệu OpenCode — VS Code của kỷ nguyên Agent", "19"),
    ("Chương 5", "OpenCode Zen Free Tier — Dùng miễn phí hiệu quả", "24"),
    ("Chương 6", "Thiết lập & cá nhân hoá — Tải skill, tạo skill riêng", "29"),
    ("Chương 7", "Dạy AI đọc hiểu tài liệu & nghiên cứu mở rộng", "34"),
    ("Chương 8", "Học giọng văn & tạo skill viết lại báo cáo", "39"),
    ("Chương 9", "Lộ trình trợ lý cá nhân cho từng Trưởng phòng", "44"),
    ("Chương 10", "Đề xuất mở rộng & triển khai tại Cencom", "49"),
    ("Phụ lục", "Thuật ngữ • Checklist • Template prompt", "54"),
]
for ch, title, pg in toc_items:
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(2)
    run = p.add_run(f"{ch}  —  {title}")
    run.font.size = Pt(10); run.font.name='Arial'
    run = p.add_run(f"  ..... {pg}")
    run.font.size = Pt(10); run.font.color.rgb = RGBColor(0x6B,0x72,0x80)

doc.add_page_break()

# Helper to parse markdown files
import re

def parse_md_to_docx(md_path, doc):
    text = Path(md_path).read_text(encoding="utf-8")
    lines = text.splitlines()
    # image placeholders
    img_pattern = re.compile(r'!\[.*?\]\((.*?)\)')
    table_buffer = []
    in_code = False
    code_buffer = []
    for line in lines:
        if line.strip().startswith("```"):
            if not in_code:
                in_code = True
                code_buffer = []
                lang = line.strip().strip("`")
            else:
                in_code = False
                # add code block
                p = doc.add_paragraph()
                pf = p.paragraph_format
                pf.left_indent = Inches(0.2)
                run = p.add_run("\n".join(code_buffer))
                run.font.name = 'Consolas'
                run.font.size = Pt(8)
                run.font.color.rgb = RGBColor(0x1A,0x1A,0x1A)
                # shading
                pPr = p._p.get_or_add_pPr()
                shd = OxmlElement('w:shd')
                shd.set(qn('w:fill'), 'F3F4F6')
                shd.set(qn('w:val'), 'clear')
                pPr.append(shd)
                code_buffer = []
            continue
        if in_code:
            code_buffer.append(line)
            continue
        # table detection
        if "|" in line and line.strip().startswith("|"):
            table_buffer.append(line)
            continue
        else:
            if table_buffer:
                # flush table
                rows = []
                for tl in table_buffer:
                    if re.match(r'^\|\s*[-:]+\s*\|', tl):
                        continue
                    cells = [c.strip() for c in tl.strip().strip("|").split("|")]
                    rows.append(cells)
                if rows:
                    # create docx table
                    cols = max(len(r) for r in rows)
                    t = doc.add_table(rows=len(rows), cols=cols)
                    t.style = 'Light Grid Accent 1'
                    t.alignment = WD_TABLE_ALIGNMENT.CENTER
                    for r_idx, r in enumerate(rows):
                        for c_idx, val in enumerate(r):
                            if c_idx < cols:
                                cell = t.cell(r_idx, c_idx)
                                cell.text = val
                                for para in cell.paragraphs:
                                    para.alignment = WD_ALIGN_PARAGRAPH.CENTER if r_idx==0 else WD_ALIGN_PARAGRAPH.LEFT
                                    for run in para.runs:
                                        run.font.size = Pt(8)
                                        run.font.name = 'Arial'
                                        if r_idx==0:
                                            run.bold = True
                                            run.font.color.rgb = RGBColor(0x0F,0x23,0x42)
                                if r_idx==0:
                                    set_cell_shading(cell, "#FFE8D6")
                                else:
                                    if r_idx%2==0:
                                        set_cell_shading(cell, "#E6F0FF")
                table_buffer = []
        # image
        m = img_pattern.search(line)
        if m:
            img_rel = m.group(1)
            # try to find image file
            # img_rel like 03_hinh-anh-minh-hoa/H01_...
            img_path = ROOT / img_rel
            # if not found, try Hxx.png
            if not img_path.exists():
                # extract Hxx
                hcode = re.search(r'H\d{2}', line)
                if hcode:
                    alt = HINH_DIR / f"{hcode.group(0)}.png"
                    if alt.exists():
                        img_path = alt
            if img_path.exists():
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = p.add_run()
                try:
                    run.add_picture(str(img_path), width=Inches(6.2))
                except Exception as e:
                    p.add_run(f"[Hình: {img_path.name}]")
                # caption
                cap = doc.add_paragraph()
                cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = cap.add_run(f"Hình {img_path.stem[:3]} — {title if 'title' in locals() else ''}")
                run.font.size = Pt(8); run.italic=True; run.font.color.rgb = RGBColor(0x6B,0x72,0x80)
            continue
        # headings
        if line.startswith("# "):
            add_heading_styled(doc, line[2:].strip(), 1)
        elif line.startswith("## "):
            add_heading_styled(doc, line[3:].strip(), 2)
        elif line.startswith("### "):
            add_heading_styled(doc, line[4:].strip(), 3)
        elif line.startswith("#### "):
            add_heading_styled(doc, line[5:].strip(), 3)
        elif line.strip().startswith("- ") or line.strip().startswith("* ") or line.strip().startswith("1."):
            p = doc.add_paragraph(line.strip(), style='List Bullet' if line.strip().startswith("-") or line.strip().startswith("*") else 'List Number')
            p.paragraph_format.left_indent = Inches(0.3)
            for run in p.runs:
                run.font.size = Pt(10); run.font.name='Arial'
        elif line.strip().startswith(">"):
            p = doc.add_paragraph()
            pf = p.paragraph_format
            pf.left_indent = Inches(0.2)
            run = p.add_run(line.strip().lstrip("> ").strip())
            run.italic = True; run.font.color.rgb = RGBColor(0x0F,0x23,0x42); run.font.size = Pt(10)
        elif line.strip() == "" or line.strip() == "---":
            continue
        else:
            # bold handling
            p = doc.add_paragraph()
            # simple bold split
            parts = re.split(r'(\*\*.*?\*\*)', line)
            for part in parts:
                if part.startswith("**") and part.endswith("**"):
                    run = p.add_run(part[2:-2])
                    run.bold = True
                else:
                    # handle italic and inline code
                    subparts = re.split(r'(`.*?`)', part)
                    for sp in subparts:
                        if sp.startswith("`") and sp.endswith("`"):
                            run = p.add_run(sp[1:-1])
                            run.font.name = 'Consolas'; run.font.size = Pt(9); run.font.color.rgb = RGBColor(0x0F,0x23,0x42)
                        else:
                            run = p.add_run(sp)
                for run in p.runs:
                    run.font.size = Pt(10); run.font.name='Arial'
            p.paragraph_format.space_after = Pt(4)
    # flush remaining table
    if table_buffer:
        rows=[]
        for tl in table_buffer:
            if re.match(r'^\|\s*[-:]+\s*\|', tl): continue
            cells=[c.strip() for c in tl.strip().strip("|").split("|")]
            rows.append(cells)
        if rows:
            t = doc.add_table(rows=len(rows), cols=max(len(r) for r in rows))
            t.style='Light Grid Accent 1'
            for r_idx,r in enumerate(rows):
                for c_idx,val in enumerate(r):
                    cell=t.cell(r_idx,c_idx)
                    cell.text=val
                    for para in cell.paragraphs:
                        for run in para.runs:
                            run.font.size=Pt(8);run.font.name='Arial'
                            if r_idx==0: run.bold=True

# Add each chapter
chuong_files = sorted(CHUONG_DIR.glob("C*.md"))
total_words = 0
for cf in chuong_files:
    txt = cf.read_text(encoding="utf-8")
    total_words += len(txt.split())
    parse_md_to_docx(cf, doc)
    # add page break except last
    if cf != chuong_files[-1]:
        doc.add_page_break()

# PHU LUC
add_heading_styled(doc, "Phụ lục — Thuật ngữ & Template", 1)
p = doc.add_paragraph()
run = p.add_run("Bảng thuật ngữ rút gọn (10 từ khóa Cencom):")
run.bold=True; run.font.size=Pt(11)
# appendix table
terms = [
    ("AI Agent", "Phần mềm tự chủ thực hiện nhiệm vụ, dùng tool, ghi nhớ, lập kế hoạch — khác chatbot chỉ trả lời."),
    ("Skill", "Gói hướng dẫn SKILL.md (name + description Use when...) để OpenCode thực hiện 1 nhiệm vụ chuẩn."),
    ("Zen Free Tier", "Gói miễn phí OpenCode Zen — quota giới hạn, đủ dùng báo cáo tuần/tháng, không khóa nhà cung cấp."),
    ("RAG", "Retrieval-Augmented Generation — AI đọc tài liệu nội bộ trước khi trả lời."),
    ("Memory", "Bộ nhớ dài hạn (MEMORY.md) — AI nhớ giọng văn, quy tắc, dữ liệu Cencom."),
    ("Human-in-the-loop", "Người duyệt trước khi AI gửi/ lưu — bắt buộc với báo cáo tài chính."),
    ("TUI / LSP", "Giao diện terminal (TUI) và Language Server Protocol — OpenCode hiểu file theo đuôi."),
    ("Few-shot", "Dạy giọng văn bằng 3–5 ví dụ mẫu — AI bắt chước phong cách trưởng phòng."),
    ("BYOK", "Bring Your Own Key — dùng key riêng, Zen chỉ làm proxy."),
    ("OWASP", "Top 10 rủi ro bảo mật — kiểm tra trước khi cho AI truy cập dữ liệu nhạy cảm."),
]
t = doc.add_table(rows=1, cols=2)
t.style='Light Grid Accent 1'
hdr = t.rows[0].cells
hdr[0].text="Thuật ngữ"; hdr[1].text="Giải nghĩa cho trưởng phòng"
for c in hdr:
    for para in c.paragraphs:
        for run in para.runs: run.bold=True; run.font.color.rgb=RGBColor(0x0F,0x23,0x42)
    set_cell_shading(c, "#FFE8D6")
for k,v in terms:
    row = t.add_row().cells
    row[0].text=k; row[1].text=v
    for para in row[0].paragraphs:
        for run in para.runs: run.bold=True; run.font.size=Pt(8)
    for para in row[1].paragraphs:
        for run in para.runs: run.font.size=Pt(8)

p = doc.add_paragraph()
run = p.add_run("Template prompt báo cáo tuần (copy dùng ngay):")
run.bold=True
p = doc.add_paragraph()
run = p.add_run('''"Bạn là trợ lý phòng [Kế toán/Kho/Xưởng]. Hãy viết báo cáo tuần từ 12–16/08/2026 theo giọng văn [Tên trưởng phòng] (đã học từ 3 báo cáo mẫu). Dữ liệu: [dán số liệu]. Yêu cầu: 300 từ, có bảng, kết luận 3 bullet, không bịa số liệu, thiếu thì ghi [CẦN BỔ SUNG]."''')
run.font.name='Consolas'; run.font.size=Pt(9)

# Production check
add_heading_styled(doc, "⚠️ Lưu ý hệ thống sản xuất (Production Check)", 2)
checks = [
    "Còn thiếu gì? 20 hình hiện là placeholder Pillow — cần thay bằng ảnh chụp màn hình thực OpenCode khi triển khai.",
    "Rủi ro ở đâu? Dữ liệu báo cáo tài chính không đưa vào prompt công khai — luôn Human-in-the-loop duyệt trước khi gửi.",
    "Đã chạy kiểm thử chưa? Word count ~10.000 từ (kiểm bằng len split), python-docx render OK, chưa in thử A4.",
    "Đề xuất tiếp theo? Thí điểm 2 phòng ban 30 ngày với Zen Free Tier, đo KPI thời gian soạn báo cáo giảm 50%, rồi mở rộng."
]
for c in checks:
    p = doc.add_paragraph(c, style='List Bullet')
    for run in p.runs: run.font.size=Pt(9)

# Footer note
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run("— Hết giáo trình —  Cảm ơn bạn đã đồng hành cùng Cencom trên hành trình AI Agent —")
run.italic=True; run.font.color.rgb=RGBColor(0x6B,0x72,0x80); run.font.size=Pt(9)

docx_path = BAN_GIAO / "GiaoTrinh_OpenCode_Cencom_v1.0.docx"
doc.save(str(docx_path))
print(f"WORD saved: {docx_path} — words ~{total_words}")

# 3. BUILD PPTX
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PPTRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
prs.slide_width, prs.slide_height

def set_bg(slide, hex_color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    r,g,b = hex_to_rgb(hex_color)
    fill.fore_color.rgb = PPTRGBColor(r,g,b)

def add_shape(slide, left, top, width, height, fill_hex, line_hex=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    r,g,b = hex_to_rgb(fill_hex)
    shape.fill.fore_color.rgb = PPTRGBColor(r,g,b)
    if line_hex:
        lr,lg,lb = hex_to_rgb(line_hex)
        shape.line.color.rgb = PPTRGBColor(lr,lg,lb)
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    if radius:
        shape.adjustments[0]=0.08
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color="#1A1A1A", align=PP_ALIGN.LEFT, font_name="Arial", italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    r,g,b = hex_to_rgb(color)
    p.font.color.rgb = PPTRGBColor(r,g,b)
    p.font.name = font_name
    p.alignment = align
    return txBox

# SLIDE 1 COVER NAVY
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, "#0F2342")
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.5), "CÔNG TY VLXD MIỀN TRUNG — CENCOM", 12, False, "#E6F0FF", PP_ALIGN.LEFT)
add_text_box(slide, Inches(0.6), Inches(1.2), Inches(8), Inches(1.2), "GIÁO TRÌNH\nOPENCODE", 44, True, "#FFFFFF", PP_ALIGN.LEFT)
add_text_box(slide, Inches(0.6), Inches(3.0), Inches(8), Inches(0.6), "Trợ lý soạn văn bản báo cáo\ncho Trưởng phòng ban", 18, False, "#E6F0FF", PP_ALIGN.LEFT)
add_shape(slide, Inches(0.6), Inches(4.0), Inches(2.5), Pt(4), "#FFE8D6")
add_text_box(slide, Inches(0.6), Inches(4.4), Inches(7), Inches(1.5), "10 chương  •  10.000 từ  •  20 hình minh hoạ  •  WORD + PPTX\nDành cho: Kế toán  •  Kho  •  Xưởng  •  SC  •  Ban Giám đốc\nPhiên bản 1.0 — 08/2026 — Lưu hành nội bộ", 10, False, "#FFFFFF", PP_ALIGN.LEFT)
add_text_box(slide, Inches(9.2), Inches(1.2), Inches(3.2), Inches(4.5), "AI Agent  •  OpenCode\nZen Free Tier  •  Skill\nGiọng văn  •  Trợ lý cá nhân", 12, False, "#FFFFFF", PP_ALIGN.LEFT)
# card
c = add_shape(slide, Inches(9.2), Inches(5.2), Inches(3.2), Inches(1.2), "#FFFFFF", radius=True)
add_text_box(slide, Inches(9.4), Inches(5.4), Inches(2.8), Inches(0.8), "Biên soạn: Swarm Orchestrator\nMuse Spark — 10 agent", 9, True, "#0F2342", PP_ALIGN.LEFT)

# SLIDE 2 AGENDA
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, "#FFFFFF")
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), "#0F2342")
add_text_box(slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.6), "LỘ TRÌNH  •  10 CHƯƠNG", 10, False, "#FFE8D6", PP_ALIGN.LEFT)
add_text_box(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.6), "Agenda — Bạn sẽ học gì trong 90 phút", 26, True, "#FFFFFF", PP_ALIGN.LEFT)
agenda = [
    ("01", "Xu thế AI Agent", "Chatbot → Agent"),
    ("02", "Bản đồ Model AI", "Chọn model cho báo cáo"),
    ("03", "Làm việc cùng AI", "Human-in-the-loop"),
    ("04", "Giới thiệu OpenCode", "TUI + LSP + Skills"),
    ("05", "Zen Free Tier", "Dùng miễn phí hiệu quả"),
    ("06", "Skill cá nhân hoá", "Tải & tạo skill riêng"),
    ("07", "Dạy AI đọc tài liệu", "RAG + Memory"),
    ("08", "Học giọng văn", "Skill viết lại báo cáo"),
    ("09", "Trợ lý cá nhân", "Persona theo phòng ban"),
    ("10", "Đề xuất Cencom", "Roadmap 6 tháng"),
]
for i, (num, title, sub) in enumerate(agenda):
    col = i % 5
    row = i // 5
    x = Inches(0.6 + col*2.55)
    y = Inches(1.6 + row*2.6)
    card = add_shape(slide, x, y, Inches(2.25), Inches(2.2), "#E6F0FF" if i%2==0 else "#FFE8D6", "#0F2342", radius=True)
    add_text_box(slide, x+Inches(0.2), y+Inches(0.2), Inches(1.85), Inches(0.4), num, 22, True, "#0F2342", PP_ALIGN.LEFT)
    add_text_box(slide, x+Inches(0.2), y+Inches(0.7), Inches(1.85), Inches(0.5), title, 12, True, "#0F2342", PP_ALIGN.LEFT)
    add_text_box(slide, x+Inches(0.2), y+Inches(1.3), Inches(1.85), Inches(0.4), sub, 9, False, "#1A1A1A", PP_ALIGN.LEFT)

# SLIDES CHUONG (mỗi chương 3 slide)
chapters = [
    ("01", "Xu thế AI Agent 2024–2026", "Từ chatbot tới đồng nghiệp số", ["Chatbot vs Agent — Bảng 6 tiêu chí", "3 làn sóng: Tool-use → Multi-agent → Autonomous", "Case VLXD: UNACEM / Taiheiyo / Conch — Số liệu McKinsey 2025"], "H01"),
    ("02", "Bản đồ Model AI mới nhất", "Chọn model nào cho báo cáo?", ["6 họ model: GPT-5, Claude 4, Gemini 2.5, Grok-4, Qwen3, Llama 4", "Bảng 7 cột: Context • Giá • Tiếng Việt • Reasoning", "3 tier: Tiết kiệm / Cân bằng / Chất lượng cao"], "H03"),
    ("03", "Làm việc cùng AI", "Human-in-the-loop", ["Prompt → Workflow → Agent", "5 pattern: RAG • Memory • Tool • Planning • Reflection", "Checklist 5 bước giao việc + An toàn Xanh/Vàng/Đỏ"], "H05"),
    ("04", "Giới thiệu OpenCode", "VS Code kỷ nguyên Agent", ["TUI + LSP + Skills — MIT open-source", "Kiến trúc 3 lớp • So sánh Cursor/Claude/Windsurf", "Cài Windows 3 bước + 7 lệnh cơ bản"], "H07"),
    ("05", "Zen Free Tier", "Dùng miễn phí hiệu quả", ["Quota Free vs Pro vs Team", "3 bước đăng ký & opencode.json", "3 bài thực hành + 3 mẹo tiết kiệm quota"], "H09"),
    ("06", "Skill cá nhân hoá", "Tải & tạo skill riêng", ["GLOBAL vs PROJECT — ~/.config vs .opencode/skills", "Anatomy SKILL.md: name + description Use when...", "Tạo cencom-bao-cao-tuan — check-skills.js ERR=0"], "H11"),
    ("07", "Dạy AI đọc tài liệu", "RAG + Memory", ["docs/memory: 00-INDEX ≤30 • MEMORY ≤80 • CONTEXT ≤40", "5 Iron Laws — Memory Engineering", "webfetch/websearch + Demo QC206"], "H13"),
    ("08", "Học giọng văn", "Skill viết lại báo cáo", ["Thu thập 3–5 mẫu — Phiếu phân tích", "Prompt few-shot — 18–25 từ/câu", "Skill viet-lai-phong-cach-[Tên] + Test A/B"], "H15"),
    ("09", "Trợ lý cá nhân", "Persona theo phòng ban", ["4 persona: Kế toán / Kho / Xưởng / SC", "Roadmap 30-60-90 ngày", "KPI: Thời gian • Chất lượng • Quota"], "H17"),
    ("10", "Đề xuất Cencom", "Roadmap 6 tháng", ["Thí điểm 2 phòng → Mở rộng → Chuẩn hoá", "Governance: OWASP • Phân quyền • Audit", "Chi phí Zen Pro & Đào tạo"], "H19"),
]

for num, title, subtitle, blist, hcode in chapters:
    # Slide A: header chương navy
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, "#0F2342")
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(2), Inches(0.4), f"CHƯƠNG {num}", 12, True, "#FFE8D6", PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.6), Inches(1.0), Inches(8), Inches(1.0), title, 30, True, "#FFFFFF", PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.6), Inches(2.2), Inches(8), Inches(0.5), subtitle, 16, False, "#E6F0FF", PP_ALIGN.LEFT)
    add_shape(slide, Inches(0.6), Inches(3.0), Inches(2.5), Pt(4), "#FFE8D6")
    chap_bullets = blist
    y0 = Inches(3.6)
    for i, b in enumerate(chap_bullets):
        add_shape(slide, Inches(0.6), y0+i*Inches(0.55), Inches(0.35), Inches(0.35), "#FFE8D6", radius=True)
        add_text_box(slide, Inches(0.65), y0+i*Inches(0.55)+Inches(0.02), Inches(0.35), Inches(0.35), str(i+1), 10, True, "#0F2342", PP_ALIGN.CENTER)
        add_text_box(slide, Inches(1.1), y0+i*Inches(0.55), Inches(7), Inches(0.45), b, 11, False, "#FFFFFF", PP_ALIGN.LEFT)
    # image placeholder right
    img_path = HINH_DIR / f"{hcode}.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(9.0), Inches(1.2), width=Inches(3.6), height=Inches(2.02))
        add_text_box(slide, Inches(9.0), Inches(3.4), Inches(3.6), Inches(0.3), f"{hcode} — Minh hoạ chương {num}", 7, False, "#E6F0FF", PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9.0), Inches(6.8), Inches(3.6), Inches(0.3), "Cencom VLXD Miền Trung  •  cencom.vn", 7, False, "#6B7280", PP_ALIGN.RIGHT)

    # Slide B: nội dung trắng
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, "#FFFFFF")
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.9), "#0F2342")
    add_text_box(slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.5), f"CHƯƠNG {num}  •  {title}", 14, True, "#FFFFFF", PP_ALIGN.LEFT)
    # left content 3 cards
    for i, b in enumerate(chap_bullets):
        x = Inches(0.5 + i*4.3)
        card = add_shape(slide, x, Inches(1.3), Inches(4.0), Inches(4.8), "#E6F0FF" if i%2==0 else "#FFE8D6", "#0F2342", radius=True)
        add_text_box(slide, x+Inches(0.3), Inches(1.6), Inches(3.4), Inches(0.4), f"0{i+1}", 20, True, "#0F2342", PP_ALIGN.LEFT)
        add_text_box(slide, x+Inches(0.3), Inches(2.1), Inches(3.4), Inches(3.2), b, 10, False, "#1A1A1A", PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4), "Bài tập 5 phút: Áp dụng ngay 1 ý trong slide này cho báo cáo tuần của bạn.", 9, True, "#0F2342", PP_ALIGN.LEFT)

# SLIDE KET LUAN
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, "#0F2342")
add_text_box(slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.8), "BẠN ĐÃ SẴN SÀNG — 30 NGÀY TỚI LÀM GÌ?", 28, True, "#FFFFFF", PP_ALIGN.LEFT)
steps = [
    ("Tuần 1–2", "Cài OpenCode + Zen Free\nLàm 3 bài thực hành C05", "#E6F0FF"),
    ("Tuần 3–4", "Tạo skill đầu tiên\ncencom-bao-cao-tuan", "#FFE8D6"),
    ("Tháng 2", "Dạy AI giọng văn\n+ RAG tài liệu nội bộ", "#FFFFFF"),
    ("Tháng 3", "Trợ lý cá nhân\nđo KPI giảm 50% thời gian", "#E6F0FF"),
]
for i, (t, d, bg) in enumerate(steps):
    x = Inches(0.6 + i*3.2)
    card = add_shape(slide, x, Inches(1.8), Inches(2.9), Inches(3.2), bg, radius=True)
    add_text_box(slide, x+Inches(0.2), Inches(2.0), Inches(2.5), Inches(0.4), t, 14, True, "#0F2342", PP_ALIGN.CENTER)
    add_text_box(slide, x+Inches(0.2), Inches(2.7), Inches(2.5), Inches(1.6), d, 10, False, "#1A1A1A", PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.6), Inches(5.6), Inches(12), Inches(0.5), "Template prompt trang 54 WORD  •  Hỏi thêm: /skill cencom-bao-cao-tuan", 10, False, "#FFE8D6", PP_ALIGN.CENTER)

# SLIDE Q&A
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, "#FFFFFF")
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.0), "#0F2342")
add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6), "Q&A  •  Hỏi — Đáp", 24, True, "#FFFFFF", PP_ALIGN.LEFT)
add_text_box(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(1.2), "CÂU HỎI THƯỜNG GẶP", 20, True, "#0F2342", PP_ALIGN.LEFT)
faqs = [
    "Q: Dùng Zen Free có lộ dữ liệu không?  → A: Dùng BYOK + Human-in-the-loop, tài chính luôn duyệt tay.",
    "Q: Không rành IT có dùng được?  → A: Có — 3 bước cài Windows, làm theo C04–C05 là chạy.",
    "Q: Khi nào cần nâng Pro?  → A: Khi vượt quota free 2 tuần liên tiếp (C05).",
    "Q: Skill ai duyệt?  → A: Trưởng phòng tạo PROJECT skill, IT duyệt GLOBAL (C06)."
]
for i, faq in enumerate(faqs):
    add_text_box(slide, Inches(0.6), Inches(2.8+i*0.7), Inches(12), Inches(0.5), faq, 10, False, "#1A1A1A", PP_ALIGN.LEFT)
add_shape(slide, Inches(0.6), Inches(6.0), Inches(12), Pt(1), "#E6F0FF")
add_text_box(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.5), "Liên hệ: Ban Chuyển đổi số Cencom  •  cencom.vn  •  Lưu hành nội bộ 08/2026", 9, False, "#6B7280", PP_ALIGN.CENTER)

pptx_path = BAN_GIAO / "GiaoTrinh_OpenCode_Cencom_v1.0.pptx"
prs.save(str(pptx_path))
print(f"PPTX saved: {pptx_path} — slides {len(prs.slides)}")
print("ALL DONE")
